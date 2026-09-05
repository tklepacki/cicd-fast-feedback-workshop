"""Builds the workshop deck from the existing one, cloning its slides as templates."""
import copy
from pptx import Presentation
from pptx.util import Pt

A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
TPL_AGENDA = 3
TPL_THEORY = 32
TPL_TASK = 36


def clone(prs, index):
    src = prs.slides[index - 1]
    dst = prs.slides.add_slide(src.slide_layout)
    for shape in list(dst.shapes):
        shape._element.getparent().remove(shape._element)
    for shape in src.shapes:
        dst.shapes._spTree.append(copy.deepcopy(shape._element))
    return dst


def text_shapes(slide):
    out = [sh for sh in slide.shapes
           if sh.has_text_frame and sh.width and sh.width > 1000000]
    out.sort(key=lambda s: s.top)
    return out


def set_text(shape, lines, size=None):
    """Replaces text, reusing the formatting of the first paragraph and run."""
    tf = shape.text_frame
    first_p = tf.paragraphs[0]
    ppr = first_p._p.find(A + 'pPr')
    ppr_tpl = copy.deepcopy(ppr) if ppr is not None else None
    rpr_tpl = None
    if first_p.runs:
        rpr = first_p.runs[0]._r.find(A + 'rPr')
        if rpr is not None:
            rpr_tpl = copy.deepcopy(rpr)

    txbody = tf._txBody
    for p in txbody.findall(A + 'p'):
        txbody.remove(p)

    for line in lines:
        text, bullet = line if isinstance(line, tuple) else (line, False)
        p = txbody.makeelement(A + 'p', {})
        if ppr_tpl is not None:
            p.append(copy.deepcopy(ppr_tpl))
        r = txbody.makeelement(A + 'r', {})
        if rpr_tpl is not None:
            r.append(copy.deepcopy(rpr_tpl))
        t = txbody.makeelement(A + 't', {})
        t.text = ('•  ' + text) if bullet else text
        r.append(t)
        p.append(r)
        txbody.append(p)

    if size:
        for p in tf.paragraphs:
            for run in p.runs:
                run.font.size = Pt(size)


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text
