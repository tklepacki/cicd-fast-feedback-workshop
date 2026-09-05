# -*- coding: utf-8 -*-
import sys, copy
sys.path.insert(0, '/tmp/deck')
from pptx import Presentation
from builder import clone, text_shapes, set_text, set_notes, TPL_AGENDA, TPL_THEORY, TPL_TASK
from content_intro import TITLE, SUBTITLE, NOTES_TITLE, NOTES_ABOUT, AGENDA, INTRO
from content_chain1 import CHAIN1
from content_chain2 import CHAIN2

prs = Presentation('/tmp/deck/base.pptx')
ORIGINAL = len(prs.slides)

# --- slajd 1: tytuł i data ---
s1 = prs.slides[0]
shapes = text_shapes(s1)
set_text(shapes[0], [TITLE])
set_text(shapes[1], [SUBTITLE + "                                   6.09.2026"])
set_notes(s1, NOTES_TITLE)
set_notes(prs.slides[1], NOTES_ABOUT)

def add(tpl, title, lines, notes):
    sl = clone(prs, tpl)
    sh = text_shapes(sl)
    body = max(sh, key=lambda s: s.height)
    head = min(sh, key=lambda s: s.top)
    if head is body and len(sh) > 1:
        head = sh[0]; body = sh[-1]
    set_text(head, [title])
    set_text(body, lines, size=18)
    set_notes(sl, notes)
    return sl

for title, lines, notes in AGENDA:
    add(TPL_AGENDA, title, lines, notes)
for title, lines, notes in INTRO:
    add(TPL_THEORY, title, lines, notes)
for kind, title, lines, notes in CHAIN1 + CHAIN2:
    add(TPL_THEORY if kind == 'theory' else TPL_TASK, title, lines, notes)

# --- usuń oryginalne slajdy 3..42, zachowując 1 i 2 ---
xml_slides = prs.slides._sldIdLst
ids = list(xml_slides)
for sid in ids[2:ORIGINAL]:
    rId = sid.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    xml_slides.remove(sid)

out = '/Users/tomaszklepacki/Desktop/workshop/docs/Szybki feedback w CICD - Warsztaty.pptx'
prs.save(out)
print(f'zapisano: {out}')
print(f'slajdów: {len(prs.slides)}')
